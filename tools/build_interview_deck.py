"""Build the final FunRec interview deck from frozen project evidence."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("docs/presentation/FunRec_Interview_Story_Final.pptx")
NAVY = (12, 20, 38)
PANEL = (23, 35, 57)
TEXT = (235, 241, 248)
MUTED = (150, 170, 194)
MINT = (41, 197, 164)

SLIDES = [
    ("FunRec：把推荐系统做成可相信的故事", "从 PyTorch 迁移、时序协议重建，到受控优化与停止规则", ["MovieLens 1M · YouTubeDNN 召回 + DeepFM 排序 + 多样性重排", "核心交付不是‘挑一个最高分模型’，而是可复现、可解释、可停止的决策链", "模型选择始终只用 Validation；封存 Test 不参与后续方向选择"], "项目代码、实验注册表与冻结结果工件"),
    ("端到端系统已迁移到 PyTorch", "离线训练与线上服务使用同一套受控工件", ["召回：YouTubeDNN 用户塔对 3,883 部电影做 exact full softmax", "排序：DeepFM 的一阶、FM 二阶与 DNN 分支共享二阶 embedding", "服务链路：多路召回 → DeepFM 排序 → 类型/年代多样性重排", "已完成离线产物、部署目录、Docker 容器内核心工件哈希一致性核验"], "docs/baseline_walkthrough.md；docs/experiment_registry.md"),
    ("先修实验协议，才讨论算法提升", "原先的评估不能支撑可信结论，因此先重建时序边界", ["每位用户按时间切分：早期交互为 Train，倒数第二条为 Validation，最后一条为 Test", "修复跨 split 用户—电影冲突与未来正样本被误采为早期负样本的问题", "Retrieval：Train 982,089；Validation 6,040；Test 6,040", "Test 只完成一次基线核验，随后封存；后续实验禁止访问、反序列化或重评估"], "docs/baseline_walkthrough.md；pytorch-baseline-v1"),
    ("可信 Baseline：真实，但有明确边界", "History-10 masked mean；Validation 选 Epoch 14", ["唯一封存 Test：Recall@10 0.137086，NDCG@10 0.070298", "Retrieval Test Recall@10 是 popularity 对照的 7.08×", "DeepFM Test ROC-AUC 0.841382；这不是线上 A/B，也不等同于 Top-K 收益", "Baseline 作为可信锚点，而非继续追逐单个总分的理由"], "pytorch-baseline-v1 tag；docs/results/baseline_v0_*.json"),
    ("Validation 诊断：整体平均数掩盖了结构性弱点", "诊断是相关证据，不直接宣称兴趣漂移因果", ["Baseline Validation Recall@10：AQ0 0.198423，而 AQ3 仅 0.105123", "AQ0−AQ3 gap = 0.093300；尾部目标 PQ0 Recall@10 = 0.079396", "高活跃用户的退化在控制目标电影热度后仍存在，但与热度有交互", "由此定义五项联合准入门槛：总 Recall、NDCG、AQ3、gap、尾部"], "docs/results/baseline_v1_validation_diagnostics.json"),
    ("门槛先冻结：不让结果反过来定义成功", "每个候选必须同时通过 5/5 项 Validation 门槛", ["Recall@10 ≥ 0.161954；NDCG@10 ≥ 0.072737；AQ3 ≥ 0.110123", "AQ0−AQ3 gap ≤ 0.088300；Tail PQ0 ≥ 0.069396", "门槛来自 Baseline 加/减预注册容忍带：既要求总体不退化，也保护弱切片", "只用 Validation 决定接受或拒绝；不因失败后调低门槛"], "docs/results/attention20_preregistered_config.json"),
    ("失败对照一：History-20 masked mean", "延长历史窗口不是免费的表达能力提升", ["唯一核心变化：History-10 → History-20，其余训练协议冻结", "Validation Recall@10 0.143543，NDCG@10 0.072490，仅通过 1/5", "被拒绝并保留；没有因为训练更复杂或窗口更长而获得晋级", "Test 始终关闭"], "docs/results/history20_masked_mean_validation_results.json"),
    ("遮罩诊断否定了‘简单历史稀释’解释", "不应把旧历史直接删掉，也不该把失败归因成单一原因", ["full-20 / recent-10 / older-10 Recall@10：0.143543 / 0.122682 / 0.071192", "只有 1/4 项诊断判据支持简单等权稀释，因此证据不足以删除旧历史", "下一步仅允许测试更有表达力的聚合机制，而不是叠加多个改动"], "docs/results/history20_mask_diagnostics.json"),
    ("Attention-20：总体显著改善，但仍按规则拒绝", "History-20 + personalized attention；Epoch 20 最佳，Epoch 23 早停", ["Recall@10 0.199007；NDCG@10 0.101495；AQ3 0.139721；Tail PQ0 0.120735", "但 AQ0−AQ3 gap = 0.120463，高于冻结上限 0.088300", "因此是 4/5 的 validation_rejected，而不是可部署的胜利模型", "这正是联合门槛的价值：阻止‘总体更好’掩盖公平性退化"], "docs/results/attention20_validation_results.json"),
    ("双时间尺度 Attention 也没有跨过决定性门槛", "保留 older-15 / recent-5，并学习融合；仍只用 Validation", ["gap 从 Attention-20 的 0.120463 降至 0.106017，但仍高于 0.088300 上限", "其余四项通过，仍是 4/5 拒绝", "同一 attention/pooling 假设族两次都在 gap 门槛失败，因此按治理规则关闭该方向"], "docs/results/dualtimescale_attention20_validation_results.json"),
    ("协议事故也进入故事：发现、隔离、修复", "透明记录失败比掩盖失败更重要", ["首个 activity-balanced 正式目录误载入含嵌入 Test 的旧工件，因反序列化 Test 而协议无效", "没有输出 Test 指标；该目录、日志和 checkpoint 全部隔离，不用于选择、报告或 PPT 指标", "修复：选择集预处理不再生成 Test；训练入口 fail-closed，只接受 train + validation", "修复代码经过 19 项相关测试，并以新目录、冻结配置重新开始"], "protocol incident；commits b4d883d/eec16bb"),
    ("失败对照二：activity-balanced loss 有效运行但方向错误", "唯一变量：仅依赖 Train 活跃度分位的 full-softmax 损失加权", ["有效替代运行：Epoch 25 早停；Validation loss 最佳 checkpoint 为 Epoch 22", "Recall@10 0.152980；NDCG@10 0.076987；AQ3 0.085828", "gap 0.127706；Tail PQ0 0.049213；仅 NDCG 通过，最终 1/5 拒绝", "说明此重加权没有改善目标切片，反而扩大高低活跃差距；Test 未访问"], "activity_balanced_history10_rerun_validation_results.json"),
    ("探索终止：把‘不继续’变成预先定义的决策", "避免在同一方向无限调参", ["同一假设族最多两个正式候选；若两次在同一决定性门槛失败则关闭该族", "若两个独立假设族均未产生 5/5 候选，则停止模型搜索、冻结可信证据链", "attention/pooling 已关闭；Train-only activity reweighting 的首个有效候选方向性失败", "因此当前决策：停止搜索，不把候选 2 当成必须完成的配额"], "docs/experiment_registry.md（前瞻性探索治理规则）"),
    ("最终推荐：部署可信 Baseline，而非伪优化候选", "用已经通过完整协议的模型交付，保留失败证据指导下一轮数据/产品工作", ["部署/演示锚点：冻结的 Baseline V1 与已核验的线上同工件链路", "不部署 Attention-20 或 activity-balanced：它们都没有通过预注册的联合准入", "后续优先级转向数据与产品验证：更多时间窗口、在线反馈、真正的 A/B 设计", "任何新模型方向需重新立项、预注册，并保留封存 Test 边界"], "实验注册表；所有候选均以 Validation 结论冻结"),
    ("面试要点：我交付的是一条可审计的决策链", "迁移能力 + 实验设计 + 对失败的诚实处理", ["工程：将训练、导出、部署、服务和验证串成可复现的 PyTorch 系统", "科学：先修数据协议，冻结门槛，按切片而非平均数做模型选择", "判断：总体指标变好仍可拒绝；发现协议事故后隔离并修复防线", "结果：留下可信 Baseline、可解释的失败轨迹，以及清晰的停止条件"], "FunRec 项目最终交付"),
]


def add_text(slide, text, left, top, width, height, size, color, bold=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor(*color)
    paragraph.font.name = "Aptos"
    return box


def build():
    if OUT.exists():
        raise FileExistsError(f"Refusing to overwrite: {OUT}")
    OUT.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width, presentation.slide_height = Inches(13.333), Inches(7.5)
    for number, (title, subtitle, bullets, source) in enumerate(SLIDES, 1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.background.fill
        background.solid(); background.fore_color.rgb = RGBColor(*NAVY)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.55), Inches(.55), Inches(.12), Inches(1.0))
        accent.fill.solid(); accent.fill.fore_color.rgb = RGBColor(*MINT); accent.line.fill.background()
        add_text(slide, title, .9, .5, 11.9, .65, 27, TEXT, True)
        add_text(slide, subtitle, .9, 1.25, 11.5, .45, 14, MUTED)
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.85), Inches(2.0), Inches(11.65), Inches(4.35))
        panel.fill.solid(); panel.fill.fore_color.rgb = RGBColor(*PANEL); panel.line.color.rgb = RGBColor(44, 63, 91)
        box = slide.shapes.add_textbox(Inches(1.2), Inches(2.35), Inches(10.95), Inches(3.65))
        frame = box.text_frame; frame.clear()
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = bullet; paragraph.font.size = Pt(18); paragraph.font.color.rgb = RGBColor(*TEXT); paragraph.font.name = "Aptos"; paragraph.space_after = Pt(15); paragraph.bullet = True
        add_text(slide, f"FUNREC / {number:02d}", .9, 6.83, 1.4, .25, 9, MINT, True)
        footer = add_text(slide, f"Evidence: {source}", 2.0, 6.83, 10.4, .25, 8, MUTED)
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    presentation.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
